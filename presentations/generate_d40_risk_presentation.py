#!/usr/bin/env python3
"""Johnson Electric Risk Committee deck — Exxsol D40 supply chain integrity."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

from generate_je_assets import ASSETS, generate_all as generate_brand

# Johnson Electric brand
ORANGE = RGBColor(0xF5, 0x82, 0x20)
ORANGE_SOFT = RGBColor(0xFD, 0xE8, 0xD0)
BLACK = RGBColor(0x14, 0x14, 0x14)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
MID_GRAY = RGBColor(0x8A, 0x8A, 0x8A)
LIGHT = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1B, 0x4F, 0x9C)
RED = RGBColor(0xC4, 0x39, 0x1D)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

SW = 13.333
SH = 7.5
OUTPUT = Path(__file__).parent / "Exxsol_D40_Supply_Chain_Risk_Committee.pptx"


def set_run(run, size=22, bold=False, color=BLACK, name="Calibri"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def add_logo(slide, left=10.85, top=6.85, width=2.15):
    path = ASSETS / "je_logo.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width))


def add_confidential(slide):
    box = slide.shapes.add_textbox(Inches(4.4), Inches(7.12), Inches(4.5), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Confidential"
    set_run(run, size=11, color=MID_GRAY)


def add_sunburst(slide):
    path = ASSETS / "je_sunburst.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(9.55), Inches(-0.4), Inches(4.2), Inches(8.3))


def add_title_bar(slide, title, subtitle=None):
    """Template 2: orange rounded bar + black title."""
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.32), Inches(0.55), Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    _no_line(bar)

    box = slide.shapes.add_textbox(Inches(1.15), Inches(0.18), Inches(11.5), Inches(0.5))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, size=30, bold=True, color=BLACK)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1.15), Inches(0.66), Inches(11.5), Inches(0.32))
        srun = sub.text_frame.paragraphs[0].add_run()
        srun.text = subtitle
        set_run(srun, size=16, color=GRAY)


def add_chrome(slide, title=None, subtitle=None, sunburst=False, logo_left=False):
    if sunburst:
        add_sunburst(slide)
    if title:
        add_title_bar(slide, title, subtitle)
    add_confidential(slide)
    if logo_left:
        add_logo(slide, left=0.4, top=6.85, width=2.15)
    else:
        add_logo(slide)


def add_orange_bullets(slide, items, left=0.7, top=1.2, width=11.8, size=24):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        set_run(run, size=size, color=BLACK)
        p.space_after = Pt(16)
        p.space_before = Pt(4)
        pPr = p._p.get_or_add_pPr()
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", "F58220")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Calibri")
        buSz = etree.SubElement(pPr, qn("a:buSzPts"))
        buSz.set("val", "2800")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "●")


def add_callout(slide, text, left=0.7, top=6.15, width=9.8, fill=ORANGE_SOFT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.58)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size=16, bold=True, color=BLACK)


def add_table(slide, headers, rows, left=0.7, top=1.25, width=12.0, col_widths=None, font=16):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(0.48 * (len(rows) + 1))
    )
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run(run, size=font, bold=True, color=WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, size=font - 1, color=BLACK)


def add_circle(slide, left, top, size, label, sub, fill=ORANGE):
    circ = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    circ.fill.background()
    circ.line.color.rgb = ORANGE
    circ.line.width = Pt(3.5)

    tbox = slide.shapes.add_textbox(Inches(left - 0.15), Inches(top + size * 0.28), Inches(size + 0.3), Inches(1.1))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_run(run, size=18, bold=True, color=BLACK)

    sbox = slide.shapes.add_textbox(Inches(left - 0.2), Inches(top + size + 0.08), Inches(size + 0.4), Inches(0.45))
    sp = sbox.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "●  " + sub
    set_run(sr, size=15, color=BLACK)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sunburst(slide)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.4), Inches(3.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exxsol D40\nSupply Chain Risk"
    set_run(run, size=44, bold=True, color=BLACK)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Quality  •  Traceability  •  Anti-Adulteration"
    set_run(r2, size=22, color=ORANGE)
    p2.space_before = Pt(18)

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Risk Committee Briefing  |  Powder Metallurgy Cleaning\nSeptember 2026"
    set_run(r3, size=18, color=GRAY)
    p3.space_before = Pt(16)

    add_confidential(slide)
    add_logo(slide, left=0.55, top=6.85, width=2.15)


def slide_exec(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Executive Summary")
    add_orange_bullets(
        slide,
        [
            "D40 cleans porous PM parts before plating and coating",
            "Cost pressure raises drum adulteration risk",
            "Fake fluid = fire, HSE and quality failures",
            "Defense: authorized supply + CoA + GC-MS",
        ],
        top=1.25,
        width=12.0,
        size=26,
    )
    add_callout(slide, "A cheaper drum can cost more than it saves.")


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Agenda")
    items = [
        ("01", "Why D40 specs\nmatter"),
        ("02", "Boiling range &\naromatics"),
        ("03", "Supply risk &\nadulteration"),
        ("04", "Controls &\nactions"),
    ]
    x = 0.85
    for num, label in items:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.0), Inches(2.4), Inches(2.4))
        circ.fill.background()
        circ.line.color.rgb = ORANGE
        circ.line.width = Pt(4)
        nbox = slide.shapes.add_textbox(Inches(x), Inches(2.55), Inches(2.4), Inches(0.5))
        np = nbox.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        nr = np.add_run()
        nr.text = num
        set_run(nr, size=28, bold=True, color=ORANGE)
        lbox = slide.shapes.add_textbox(Inches(x - 0.1), Inches(3.1), Inches(2.6), Inches(1.0))
        lp = lbox.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = label
        set_run(lr, size=16, bold=True, color=BLACK)
        x += 3.05

    add_callout(slide, "Focus: protect solvent integrity from plant to production line.", top=5.85, width=12.0)


def slide_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Business Context", "Why D40 is on the Risk Agenda")
    add_orange_bullets(
        slide,
        [
            "PM pores trap oils — need residue-free clean",
            "D40: fast dry, low odor, metal-safe",
            "Not a commodity — specs are critical",
            "Look-alike solvents raise fraud risk",
        ],
        size=26,
    )


def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Product Profile", "Exxsol D40 key specifications")
    add_table(
        slide,
        ["Spec", "Value", "Why it matters"],
        [
            ["Distillation (IBP–DP)", "163–187 °C", "Uniform drying; no heavy-end residue"],
            ["Flash point", "~48 °C", "Flammable — ignition controls required"],
            ["Aromatics", "<0.1%", "Low odor, safer, clean surface"],
            ["OEL", "1,200 mg/m³", "4× safer than white spirit"],
            ["KB / viscosity", "~31  /  ~1.28 mm²/s", "Controlled solvency; pore penetration"],
        ],
        col_widths=[3.2, 3.4, 5.4],
        font=17,
    )


def slide_boiling(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Boiling Range", "IBP  •  Dry Point  •  Narrow window")
    add_orange_bullets(
        slide,
        [
            "IBP = first drop; Dry Point = last drop",
            "Narrow ~24 °C range = predictable drying",
            "Wide range leaves heavy ends in PM pores",
            "MC / methanol collapses the profile",
        ],
        size=26,
    )
    add_callout(slide, "Adulteration = residue in pores and unstable drying.")


def slide_aromatic(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Low Aromatics", "Dearomatized = engineered safety")
    add_orange_bullets(
        slide,
        [
            "Less odor and toxicity for operators",
            "Safer for seals and elastomers",
            "Clean surface for plating / coating",
            "Adulteration voids the certified profile",
        ],
        size=26,
    )


def slide_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "D40 vs D60", "Same family — different process window")
    add_table(
        slide,
        ["", "Exxsol D40", "Exxsol D60 / D60(S)"],
        [
            ["Boiling range", "163–187 °C", "180–210 °C"],
            ["Flash point", "~48 °C", "~68 °C"],
            ["Dry speed", "Faster", "Slower"],
            ["Best for", "PM fast clean", "Higher flash-point needs"],
        ],
        col_widths=[3.0, 4.5, 4.5],
        font=18,
    )
    add_callout(slide, "D60(S) is a supplier designation — not a different formula.", top=5.55)


def slide_supply_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Supply Chain Risk", "Price pressure raises fraud incentive")
    add_orange_bullets(
        slide,
        [
            "Middle East conflict stresses feedstock",
            "ExxonMobil +$0.06/lb (Mar 2026)",
            "China terminal price +5% and rising",
            "Wide price gap invites adulteration",
        ],
        size=26,
    )
    add_callout(slide, "Do not buy on price alone. Verify the chain.")


def slide_adulteration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Adulteration Threat", "How unscrupulous traders cut cost")
    add_orange_bullets(
        slide,
        [
            "Mix cheaper MC and methanol into D40",
            "Reuse genuine drums and fake CoA / labels",
            "Looks clear — visual check is not enough",
            "Illegal — voids SDS and fire-risk assessment",
        ],
        size=26,
    )


def slide_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Impact of Fraud")
    add_table(
        slide,
        ["Property", "Genuine D40", "Adulterated"],
        [
            ["Distillation", "163–187 °C", "Starts ~40 °C"],
            ["Flash point", "~48 °C", "Much lower (methanol ~12 °C)"],
            ["Cleaning", "Controlled, residue-free", "Fails or attacks parts"],
            ["Downstream", "Clean active surface", "Plating / coating peel"],
        ],
        col_widths=[3.0, 4.5, 4.5],
        font=17,
    )


def slide_enterprise(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Enterprise Risk")
    items = [
        ("Safety", "Fire & SDS breach"),
        ("Quality", "Plating / coating fail"),
        ("Operations", "Scrap & downtime"),
        ("Reputation", "Customer stop-ship"),
    ]
    x = 0.7
    for title, body in items:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.85), Inches(3.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ORANGE
        card.line.width = Pt(2.25)
        t = slide.shapes.add_textbox(Inches(x + 0.12), Inches(2.2), Inches(2.6), Inches(0.8))
        tp = t.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = title
        set_run(tr, size=22, bold=True, color=ORANGE)
        b = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.15), Inches(2.55), Inches(1.3))
        bp = b.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = body
        set_run(br, size=18, color=BLACK)
        x += 3.1
    add_callout(slide, "One adulterated lot can hit safety, quality and customer delivery together.")


def slide_qc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Distributor QC", "Three lines of defense")
    add_orange_bullets(
        slide,
        [
            "Buy only ExxonMobil-authorized distributors",
            "Batch CoA before goods receipt",
            "Check IBP, flash point and aromatics",
            "Inspect seals, holograms and drum labels",
        ],
        size=26,
    )


def slide_trace(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Traceability", "Chain of custody")
    add_orange_bullets(
        slide,
        [
            "Track lot: plant → ISO tank → drum → line",
            "Barcode drum to CoA in ERP",
            "Annual distributor audit + right-to-audit",
            "Quarantine suspect lots in < 4 hours",
        ],
        size=26,
    )


def slide_detect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Detection Controls")
    add_table(
        slide,
        ["Control", "When", "Goal"],
        [
            ["CoA review", "Every batch", "Match ExxonMobil spec"],
            ["On-site screen", "Each lot", "Density / flash / odor flag"],
            ["GC-MS test", "Quarterly + anomalies", "Prove purity / ID adulterants"],
            ["Price check", "Ongoing", "Spot “too cheap” offers"],
        ],
        col_widths=[3.5, 4.0, 4.5],
        font=17,
    )


def slide_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Recommended Actions", "Risk Committee approval requested")
    actions = [
        "Authorized distributors only",
        "CoA + quarantine all D40 receipts",
        "Fund quarterly GC-MS testing",
        "ERP drum-to-line traceability",
        "Distributor audit within 90 days",
    ]
    y = 1.35
    for i, text in enumerate(actions, 1):
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.75), Inches(y), Inches(0.48), Inches(0.48))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ORANGE
        _no_line(circ)
        cp = circ.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i)
        set_run(cr, size=16, bold=True, color=WHITE)
        box = slide.shapes.add_textbox(Inches(1.45), Inches(y + 0.04), Inches(10.8), Inches(0.48))
        br = box.text_frame.paragraphs[0].add_run()
        br.text = text
        set_run(br, size=24, color=BLACK)
        y += 0.88


def slide_kpi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Governance & KPIs")
    add_orange_bullets(
        slide,
        [
            "Approve Solvent Integrity Program",
            "100% CoA verified before use",
            "Zero unauthorized suppliers",
            "Quarterly GC-MS minimum",
        ],
        size=26,
    )
    add_callout(slide, "Prevention costs less than one fire, scrap event or stop-ship.")


def slide_appendix_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sunburst(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.2), Inches(2.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Appendix"
    set_run(run, size=48, bold=True, color=BLACK)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Supply chain  •  GC-MS quality control"
    set_run(r2, size=22, color=ORANGE)
    add_confidential(slide)
    add_logo(slide, left=0.55, top=6.85, width=2.15)


def _picture_slide(prs, title, subtitle, image_name, bullets, img_left=6.55, img_top=1.15, img_w=6.3, img_h=5.4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, title, subtitle)
    path = ASSETS / image_name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(img_left), Inches(img_top), Inches(img_w), Inches(img_h))
    add_orange_bullets(slide, bullets, left=0.55, top=1.2, width=5.8, size=20)
    return slide


def slide_supply_flow(prs):
    """Plant → ISO tank → distributor repack + QC."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "D40 Supply Chain", "Plant  →  ISO tank  →  authorized distributor")
    path = ASSETS / "supply_hydrocarbon.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.5), Inches(1.1), Inches(12.4), Inches(5.15))
    add_logo(slide)
    add_confidential(slide)


def slide_supply_trace(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Traceability Controls", "Seals, holograms, QR and CoA at every handoff")
    path = ASSETS / "supply_traceability.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.5), Inches(1.1), Inches(12.4), Inches(5.15))


def slide_supply_qc_points(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Repackaging — Highest Risk Point")
    add_orange_bullets(
        slide,
        [
            "Plant fills ISO tanks with sealed, labeled product",
            "Ocean / air move under GPS and manifest control",
            "Distributor repack into drums is the fraud window",
            "Demand: tamper seals, holographic labels, batch CoA",
        ],
        size=24,
        top=1.25,
        width=12.0,
    )
    add_callout(slide, "Never accept a drum that cannot be traced back to the ISO-tank lot.")


def slide_gcms_equip(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "GC-MS Equipment", "Gold-standard identity test for D40")
    path = ASSETS / "gcms_equipment.jpg"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.45), Inches(1.15), Inches(8.15), Inches(5.15))
    add_orange_bullets(
        slide,
        [
            "GC separates the mix",
            "MS fingerprints each peak",
            "Finds MC, methanol, others",
            "Not visible by eye",
        ],
        left=8.8,
        top=1.35,
        width=4.1,
        size=18,
    )


def slide_gcms_curve(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "How GC-MS Detects Abnormalities")
    path = ASSETS / "gcms_curve.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.45), Inches(1.1), Inches(8.15), Inches(5.2))
    add_orange_bullets(
        slide,
        [
            "Genuine: even C8–C14 peaks",
            "Fake: extra / irregular peaks",
            "Contaminants = fraud signal",
            "Use as lot fingerprint",
        ],
        left=8.8,
        top=1.35,
        width=4.1,
        size=18,
    )


def slide_gcms_peak_findings(prs):
    """Appendix slide: speculative reading of labeled giants on this lot chromatogram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(
        slide,
        "GC-MS Peak Findings — This Lot",
        "Speculative IDs from RT order in a C9–C11 D40 cut  •  confirm vs method + MS library",
    )

    path = ASSETS / "gcms_lot_chromatogram.jpg"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(7.15), Inches(1.12), Inches(5.75), Inches(4.55))

    add_table(
        slide,
        ["RT (min)", "Likely family", "QC reading"],
        [
            ["3.765", "C8–C9 light end (IBP)", "Keep tiny — not MC / MeOH"],
            ["7.825", "C9 / early-C10 isomers", "First body cluster"],
            ["8.583  ★", "C10 iso / naphthene", "Giant — heart of the cut"],
            ["9.030", "C10 isomer satellite", "Should track the 8.583 giant"],
            ["9.822 / 9.954", "C10–C11 isomer pair", "Adjacent boiling points"],
            ["10.530  ★", "C11 iso / naphthene", "2nd giant — ratio vs 8.583"],
            ["13.179 / 13.545", "C11–C12 heavy end (DP)", "Keep small — pore-residue risk"],
        ],
        left=0.45,
        top=1.10,
        width=6.55,
        col_widths=[1.55, 2.35, 2.65],
        font=12,
    )
    add_callout(
        slide,
        "Read: two giants (8.583 & 10.530 min) = D40 fingerprint. Growing 3.765 = light adulterant. Growing 13.x = wrong / wide cut.",
        left=0.45,
        top=6.18,
        width=12.4,
    )
    return slide


def slide_gcms_why(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Why GC-MS Matters for QC")
    add_orange_bullets(
        slide,
        [
            "Proves the drum is Exxsol D40 — not a look-alike",
            "Identifies MC (~40 °C) and methanol (~65 °C) peaks",
            "Protects flash-point, residue-free dry and worker safety",
            "Creates auditable evidence for IATF / PPAP material control",
        ],
        size=24,
    )
    add_callout(slide, "Recommended: GC-MS every critical lot + quarterly random drums.")


def slide_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sunburst(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.2), Inches(2.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Questions?"
    set_run(run, size=48, bold=True, color=BLACK)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Risk Committee  •  Exxsol D40 Integrity"
    set_run(r2, size=22, color=ORANGE)
    add_confidential(slide)
    add_logo(slide, left=0.55, top=6.85, width=2.15)


def build_presentation():
    generate_brand()
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    slide_title(prs)
    slide_exec(prs)
    slide_agenda(prs)
    slide_context(prs)
    slide_product(prs)
    slide_boiling(prs)
    slide_aromatic(prs)
    slide_compare(prs)
    slide_supply_risk(prs)
    slide_adulteration(prs)
    slide_impact(prs)
    slide_enterprise(prs)
    slide_qc(prs)
    slide_trace(prs)
    slide_detect(prs)
    slide_actions(prs)
    slide_kpi(prs)
    slide_appendix_divider(prs)
    slide_supply_flow(prs)
    slide_supply_trace(prs)
    slide_supply_qc_points(prs)
    slide_gcms_equip(prs)
    slide_gcms_curve(prs)
    slide_gcms_peak_findings(prs)
    slide_gcms_why(prs)
    slide_qa(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides)")
    return prs


def build_appendix_insert():
    """Single-slide file to copy into the full pack appendix."""
    generate_brand()
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide_gcms_peak_findings(prs)
    out = Path(__file__).parent / "Appendix_GCMS_Peak_Findings.pptx"
    prs.save(out)
    print(f"Saved insert slide: {out}")


if __name__ == "__main__":
    build_presentation()
    build_appendix_insert()
